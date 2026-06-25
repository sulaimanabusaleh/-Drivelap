"""
create_presentation.py — Erstellt eine PowerPoint-Präsentation über das DriveLab RL-Projekt.

Starten:
    pip install python-pptx
    python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# Farben
DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT     = RGBColor(0x89, 0xB4, 0xFA)   # Blau
GREEN      = RGBColor(0xA6, 0xE3, 0xA1)   # Grün
RED        = RGBColor(0xF3, 0x8B, 0xA8)   # Rot
YELLOW     = RGBColor(0xF9, 0xE2, 0xAF)   # Gelb
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY  = RGBColor(0x31, 0x32, 0x4F)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # Blank layout


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, text, font_size=20,
        bold=False, color=WHITE, bg_color=None, align=PP_ALIGN.LEFT,
        italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def title_box(slide, text, sub=None):
    box(slide, 0, 0, 13.33, 1.1, text,
        font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    if sub:
        box(slide, 0, 1.0, 13.33, 0.6, sub,
            font_size=18, color=GRAY, align=PP_ALIGN.CENTER, italic=True)


def divider(slide, top=1.5):
    line = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(0.4), Inches(top), Inches(12.93), Inches(top)
    )
    line.line.color.rgb = ACCENT
    line.line.width = Pt(1.5)


def bullet_box(slide, left, top, width, height, items, title=None,
               title_color=ACCENT, item_color=WHITE, font_size=16,
               bg_color=DARK_GRAY, spacing=0.38):
    rect = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_color
    tf = rect.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(font_size + 2)
        run.font.bold = True
        run.font.color.rgb = title_color

    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        if first and not title:
            first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = item_color


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 1 — Titelfolie
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 1.8, 13.33, 1.5,
    "DriveLab — Spurhalten mit Reinforcement Learning",
    font_size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
box(s, 0, 3.2, 13.33, 0.8,
    "Entwicklung eines PPO-Agenten für autonomes Fahrspurhalten",
    font_size=22, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
box(s, 0, 4.1, 13.33, 0.6,
    f"Bachelorarbeit Vorbereitung  |  {datetime.date.today().strftime('%B %Y')}",
    font_size=16, color=GRAY, align=PP_ALIGN.CENTER)
box(s, 0, 5.5, 13.33, 0.5,
    "Sulaiman Abusaleh",
    font_size=18, color=GREEN, align=PP_ALIGN.CENTER, bold=True)


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 2 — Projektziel
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
title_box(s, "Projektziel", "Was soll der Agent lernen?")
divider(s)

box(s, 0.5, 1.8, 12.3, 1.0,
    "Ein KI-Agent soll ein Fahrzeug automatisch in der Fahrspurmitte halten —",
    font_size=20, color=WHITE)
box(s, 0.5, 2.5, 12.3, 0.8,
    "auf verschiedenen Straßen, ohne menschliche Steuerung.",
    font_size=20, color=YELLOW, bold=True)

bullet_box(s, 0.5, 3.3, 5.8, 3.5, [
    "  Fahrzeug bleibt in der Spur",
    "  Querversatz (e_y) minimieren",
    "  Kurswinkel (e_psi) minimieren",
    "  Kurven sicher durchfahren",
    "  4 verschiedene Strecken meistern",
], title="✓  Ziele:", bg_color=DARK_GRAY)

bullet_box(s, 6.5, 3.3, 6.3, 3.5, [
    "  Reinforcement Learning (PPO)",
    "  C++ Simulation (DriveLab)",
    "  Python 3.12 + 3.13 Architektur",
    "  Stabile Baselines 3 (SB3)",
    "  Gymnasium Interface",
], title="⚙  Technologien:", bg_color=DARK_GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 3 — System-Architektur
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
title_box(s, "System-Architektur", "Zwei-Prozess-Bridge (Python 3.12 ↔ Python 3.13)")
divider(s)

# Prozess 1
bullet_box(s, 0.3, 1.7, 5.5, 4.8, [
    "",
    "  train.py / eval.py",
    "  drivelab_env.py",
    "  gymnasium.Env",
    "  Stable-Baselines3 (PPO)",
    "  VecNormalize",
    "  Plotter / TensorBoard",
    "",
], title="  Python 3.12 (py312)", title_color=GREEN,
   bg_color=RGBColor(0x1A, 0x2F, 0x1A))

# Pfeil
box(s, 5.8, 3.5, 1.7, 0.6, "JSON\n←→\nstdin/stdout",
    font_size=13, color=YELLOW, align=PP_ALIGN.CENTER, bold=True)

# Prozess 2
bullet_box(s, 7.5, 1.7, 5.5, 4.8, [
    "",
    "  drivelab_server.py",
    "  drivelab.cp313-win_amd64.pyd",
    "  C++ Fahrzeugdynamik",
    "  Straßendefinitionen (JSON)",
    "  Physik-Simulation",
    "  track.csv Ausgabe",
    "",
], title="  Python 3.13 (base)", title_color=ACCENT,
   bg_color=RGBColor(0x1A, 0x1A, 0x2F))

box(s, 0.3, 6.3, 12.7, 0.6,
    "Grund: drivelab.pyd wurde nur für Python 3.13 kompiliert — SB3/Gymnasium benötigt Python 3.12",
    font_size=13, color=GRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 4 — Observation Space
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
title_box(s, "Observation Space", "Was sieht der Agent?  →  8 Werte pro Zeitschritt")
divider(s)

obs = [
    ("e_y",        "Querversatz zur Fahrbahnmitte",      "−5.0 … +5.0 m",      ACCENT),
    ("e_psi",      "Kurswinkel-Fehler",                  "−1.5 … +1.5 rad",    ACCENT),
    ("curvature",  "Straßenkrümmung voraus",             "−0.2 … +0.2",        GREEN),
    ("distLeft",   "Abstand zum linken Rand",            "0 … 10 m",           GREEN),
    ("distRight",  "Abstand zum rechten Rand",           "0 … 10 m",           GREEN),
    ("laneWidth",  "Spurbreite",                         "0 … 10 m",           YELLOW),
    ("v",          "Fahrzeuggeschwindigkeit",            "0 … 50 m/s",         YELLOW),
    ("psidot",     "Gierrate (Drehgeschwindigkeit)",     "−5 … +5 rad/s",      RED),
]

for i, (name, desc, rang, col) in enumerate(obs):
    row_top = 1.7 + i * 0.67
    rect = s.shapes.add_shape(1, Inches(0.4), Inches(row_top),
                               Inches(12.5), Inches(0.58))
    rect.fill.solid()
    rect.fill.fore_color.rgb = DARK_GRAY if i % 2 == 0 else RGBColor(0x28, 0x28, 0x40)
    rect.line.fill.background()
    box(s, 0.5, row_top + 0.05, 2.0, 0.5, name, font_size=15, bold=True, color=col)
    box(s, 2.5, row_top + 0.05, 6.5, 0.5, desc, font_size=14, color=WHITE)
    box(s, 9.2, row_top + 0.05, 3.5, 0.5, rang, font_size=13, color=GRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 5 — Action Space & Reward
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
title_box(s, "Action Space & Reward-Funktion", "Was tut der Agent und wie wird er bewertet?")
divider(s)

bullet_box(s, 0.4, 1.7, 5.8, 2.5, [
    "",
    "  1 Wert:  Lenkwinkel",
    "  Bereich:  −0.5 … +0.5 rad",
    "  Gas:  fix 0.0 (konstant)",
    "  Bremse:  fix 0.0 (konstant)",
    "",
], title="  Action Space (1D, kontinuierlich)", bg_color=DARK_GRAY)

bullet_box(s, 0.4, 4.3, 5.8, 2.8, [
    "",
    "  Zukünftig: 3D Action Space",
    "  [Lenkwinkel, Gas, Bremse]",
    "  Erweiterung für Bachelorarbeit",
    "",
], title="  Geplante Erweiterung", title_color=YELLOW, bg_color=RGBColor(0x2F, 0x2F, 0x1A))

bullet_box(s, 6.5, 1.7, 6.4, 2.5, [
    "",
    "  r = +1.0  (pro Schritt)",
    "  − 2.0 × e_y²  (Querversatz-Strafe)",
    "  − 0.5 × e_psi²  (Winkel-Strafe)",
    "  = −10.0  (bei Crash / offRoad)",
    "",
], title="  Reward-Funktion", bg_color=DARK_GRAY)

box(s, 6.5, 4.3, 6.4, 0.5,
    "Interpretation:", font_size=15, bold=True, color=ACCENT)
box(s, 6.5, 4.75, 6.4, 0.5,
    "  ●  Perfektes Spurhalten  →  ~+1.0 / Schritt",
    font_size=14, color=GREEN)
box(s, 6.5, 5.15, 6.4, 0.5,
    "  ●  Spurversatz 1m  →  +1 − 2(1²) = −1.0",
    font_size=14, color=YELLOW)
box(s, 6.5, 5.55, 6.4, 0.5,
    "  ●  Straßenverlassen  →  −10.0 (Episode Ende)",
    font_size=14, color=RED)


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 6 — PPO Hyperparameter
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
title_box(s, "PPO Hyperparameter", "Proximal Policy Optimization — Konfiguration")
divider(s)

params = [
    ("learning_rate",    "3 × 10⁻⁴",   "Lerngeschwindigkeit des neuronalen Netzes"),
    ("n_steps",          "2048",        "Schritte pro Umgebung vor jedem Update"),
    ("batch_size",       "64",          "Mini-Batch-Größe beim Training"),
    ("n_epochs",         "10",          "Wiederholungen pro Update"),
    ("gamma",            "0.99",        "Diskontierungsfaktor für zukünftige Rewards"),
    ("gae_lambda",       "0.95",        "Generalized Advantage Estimation"),
    ("clip_range",       "0.2",         "PPO Clipping — verhindert zu große Policy-Änderungen"),
    ("ent_coef",         "0.01",        "Entropie-Koeffizient — fördert Erkundung"),
    ("n_envs",           "4",           "Parallele Trainingsumgebungen"),
    ("total_timesteps",  "500.000+",    "Gesamte Trainingsschritte"),
]

for i, (param, val, desc) in enumerate(params):
    row_top = 1.7 + i * 0.55
    col = ACCENT if i % 2 == 0 else GREEN
    box(s, 0.5, row_top, 3.0, 0.5, param, font_size=14, bold=True, color=col)
    box(s, 3.5, row_top, 1.8, 0.5, val,   font_size=14, color=YELLOW, bold=True)
    box(s, 5.3, row_top, 7.8, 0.5, desc,  font_size=13, color=GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 7 — Trainings-Pipeline
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
title_box(s, "Trainings-Pipeline", "Wie läuft ein Training ab?")
divider(s)

steps = [
    ("1", "train.py starten",        "python train.py  →  erstellt automatisch run_1, run_2, ...",           ACCENT),
    ("2", "Parallele Umgebungen",    "4× DrivelabEnv starten  →  4 Subprocess-Server (Python 3.13)",         GREEN),
    ("3", "VecNormalize",            "Beobachtungen normalisieren  →  stabiles Training",                     YELLOW),
    ("4", "PPO lernt",               "Agent wählt Lenkwinkel, Simulation reagiert, Reward berechnet",         WHITE),
    ("5", "EvalCallback",            "Alle 10.000 Schritte: bestes Modell speichern (best_model.zip)",        GREEN),
    ("6", "CheckpointCallback",      "Alle 10.000 Schritte: Zwischenstand speichern (checkpoints/)",          ACCENT),
    ("7", "EvalAndPlotCallback",     "Nach jedem Update: Fahrt simulieren + 5 Plots speichern",               YELLOW),
    ("8", "vecnormalize.pkl",        "Nach jedem Update gespeichert  →  eval.py immer nutzbar",               RED),
    ("9", "TensorBoard",             "Live-Monitoring in Browser  →  results/<run>/tensorboard/",             ACCENT),
]

for i, (nr, title_t, desc, col) in enumerate(steps):
    top = 1.65 + i * 0.58
    box(s, 0.3, top, 0.5, 0.5, nr, font_size=16, bold=True, color=col,
        align=PP_ALIGN.CENTER)
    box(s, 0.9, top, 3.5, 0.5, title_t, font_size=14, bold=True, color=col)
    box(s, 4.4, top, 8.7, 0.5, desc, font_size=13, color=GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 8 — Ergebnis-Ordnerstruktur
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
title_box(s, "Ergebnis-Ordnerstruktur", "Was wird wo gespeichert?")
divider(s)

code = (
    "results/\n"
    "├── run_1/\n"
    "│   ├── best_model/\n"
    "│   │   └── best_model.zip     ← bestes Modell (höchster Eval-Reward)\n"
    "│   ├── checkpoints/\n"
    "│   │   └── ppo_drivelab_10000_steps.zip  ← Zwischenstände\n"
    "│   ├── plots/\n"
    "│   │   └── update_0001/       ← 5 PNG-Plots pro Update\n"
    "│   ├── tensorboard/           ← TensorBoard-Logs\n"
    "│   ├── vecnormalize.pkl       ← Normierungs-Statistiken\n"
    "│   └── eval_trajectory.csv   ← Trajektorie aus eval.py\n"
    "├── run_2/  (nächstes Training → eigener Ordner)\n"
    "└── track.csv                  ← Streckengeometrie\n"
)

bullet_box(s, 0.4, 1.7, 7.5, 5.4,
           [code], bg_color=RGBColor(0x0D, 0x0D, 0x1A))

bullet_box(s, 8.2, 1.7, 4.8, 5.4, [
    "",
    "  best_model.zip",
    "  → gespeichertes neuronales Netz",
    "  → policy_net + value_net",
    "",
    "  vecnormalize.pkl",
    "  → Mittelwert + Standardabw.",
    "  → MUSS zu Modell passen!",
    "",
    "  Jeder Run ist isoliert",
    "  → kein Überschreiben",
    "",
], title="  Wichtig:", bg_color=DARK_GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 9 — Mehrere Straßen
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
title_box(s, "Trainings-Straßen", "4 Strecken für robustes Training")
divider(s)

roads = [
    ("1", "newRoad",               "Einfache Kurve nach rechts",                 "~315 m",  GREEN),
    ("2", "einfacher_Rundkurs",    "Ovaler Rundkurs (geschlossen)",              "~700 m",  ACCENT),
    ("3", "road3_kurvig",          "S-Kurven (abwechselnd links/rechts ±0.04)",  "~675 m",  YELLOW),
    ("4", "road4_komplex",         "Lange Geraden mit variierenden Kurven",      "~666 m",  RED),
]

for i, (nr, name, desc, length, col) in enumerate(roads):
    top = 1.8 + i * 1.1
    rect = s.shapes.add_shape(1, Inches(0.4), Inches(top),
                               Inches(12.5), Inches(0.95))
    rect.fill.solid()
    rect.fill.fore_color.rgb = DARK_GRAY
    rect.line.fill.background()
    box(s, 0.6, top + 0.1, 0.5, 0.75, nr,     font_size=22, bold=True, color=col)
    box(s, 1.2, top + 0.1, 3.8, 0.4,  name,   font_size=16, bold=True, color=col)
    box(s, 1.2, top + 0.5, 8.0, 0.4,  desc,   font_size=14, color=WHITE)
    box(s, 10.5, top + 0.1, 2.3, 0.75, length, font_size=18, bold=True,
        color=YELLOW, align=PP_ALIGN.RIGHT)

box(s, 0.4, 6.25, 12.5, 0.5,
    "Training: zufällige Strecke pro Episode  |  Eval: Strecke wählbar (1–4) oder interaktives Menü",
    font_size=14, color=GRAY, italic=True)

box(s, 0.4, 6.7, 12.5, 0.5,
    "Teststrecke road5_test: Geschlossener Kurs mit Chicanes + S-Kurven (~1750 m) — für Generalisierungstest",
    font_size=13, color=YELLOW, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 10 — TensorBoard Metriken
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
title_box(s, "TensorBoard — Training überwachen", "Live-Monitoring im Browser")
divider(s)

metrics = [
    ("ep_rew_mean",         "Steigt → positiv",   "Bleibt negativ",    "Ø Reward pro Episode",              ACCENT),
    ("ep_len_mean",         "Steigt ↑",           "Fällt ↓",           "Auto bleibt länger auf Straße",     ACCENT),
    ("value_loss",          "Fällt → 0 ↓",        "Steigt wieder ↑",   "KI schätzt Zukunft besser",         GREEN),
    ("explained_variance",  "Steigt → 1.0 ↑",     "Bleibt bei 0",      "Reward-Vorhersage (0=schlecht)",    GREEN),
    ("approx_kl",           "Klein ~0.01",         "Über 0.02",         "Policy ändert sich kontrolliert",   YELLOW),
    ("entropy_loss",        "Fällt langsam ↓",     "Fällt zu schnell",  "KI erkundet noch genug",            YELLOW),
    ("train/std",           "Fällt → ~0.3",        "Steigt auf 3+",     "Policy-Ausgabe stabil/instabil",    RED),
]

box(s, 0.5, 1.65, 4.0, 0.45, "Metrik",       font_size=14, bold=True, color=ACCENT)
box(s, 4.5, 1.65, 3.0, 0.45, "Gut wenn",     font_size=14, bold=True, color=GREEN)
box(s, 7.5, 1.65, 3.0, 0.45, "Schlecht wenn",font_size=14, bold=True, color=RED)
box(s, 10.5, 1.65, 2.5, 0.45, "Bedeutung",   font_size=14, bold=True, color=GRAY)

for i, (name, gut, schlecht, bedeutung, col) in enumerate(metrics):
    top = 2.1 + i * 0.67
    bg_c = DARK_GRAY if i % 2 == 0 else RGBColor(0x28, 0x28, 0x40)
    rect = s.shapes.add_shape(1, Inches(0.4), Inches(top),
                               Inches(12.5), Inches(0.6))
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_c
    rect.line.fill.background()
    box(s, 0.5, top+0.07, 3.8, 0.5, name,       font_size=13, bold=True, color=col)
    box(s, 4.5, top+0.07, 2.8, 0.5, gut,         font_size=12, color=GREEN)
    box(s, 7.5, top+0.07, 2.8, 0.5, schlecht,    font_size=12, color=RED)
    box(s, 10.5, top+0.07, 2.5, 0.5, bedeutung,  font_size=11, color=GRAY)

box(s, 0.5, 6.8, 12.0, 0.45,
    "Starten:  tensorboard --logdir results/    →    Browser: localhost:6006",
    font_size=14, color=YELLOW, bold=True)


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 11 — Evaluation
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
title_box(s, "Evaluation — Modell testen", "eval.py und modell_test.py")
divider(s)

bullet_box(s, 0.4, 1.7, 6.0, 3.0, [
    "",
    "  python eval.py",
    "  → Interaktives Straßen-Menü",
    "",
    "  python eval.py run_5/best_model.zip 2",
    "  → Direkt Straße 2 (Rundkurs)",
    "",
    "  Zeigt 5 Plots am Ende",
    "  Speichert eval_trajectory.csv",
    "",
], title="  eval.py — Bekannte Strecken", bg_color=DARK_GRAY)

bullet_box(s, 6.7, 1.7, 6.2, 3.0, [
    "",
    "  python modell_test.py",
    "  → Lädt road5_test.json",
    "  → NEUE unbekannte Strecke",
    "",
    "  Testet Generalisierbarkeit:",
    "  Hat das Modell wirklich gelernt",
    "  oder nur auswendig gelernt?",
    "",
], title="  modell_test.py — Neue Strecke", bg_color=DARK_GRAY)

bullet_box(s, 0.4, 4.8, 12.5, 2.2, [
    "  model_path  =  results/run_5/best_model/best_model.zip",
    "  vecnormalize_path  =  results/run_5/vecnormalize.pkl     ← MUSS zum Modell passen!",
    "  road  =  Auswahl 1-4 oder automatisch (None = zufällig)",
    "  Plots: Trajektorie, Querversatz, Heading, Randabstände, Geschwindigkeit",
], title="  Wichtige Parameter:", bg_color=RGBColor(0x1A, 0x2F, 0x1A))


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 12 — Probleme & Lösungen
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
title_box(s, "Probleme & Lösungen", "Was wurde gefunden und behoben?")
divider(s)

problems = [
    ("sys.path Bug",          "drivelab.pyd wurde nicht gefunden",           "str(HERE) statt str(HERE/\"drivelab...\")",    GREEN),
    ("Obs-Shape Mismatch",    "5D vecnormalize mit 8D Modell laden",         "Neu trainieren mit korrekter 8D Obs",          GREEN),
    ("Auto fährt über Ende",  "Episode lief über Streckenende hinaus",       "ROAD_LENGTH berechnen + terminieren",          GREEN),
    ("vecnormalize fehlt",    "Training unterbrochen → eval.py crasht",     "pkl nach JEDEM Update speichern",              GREEN),
    ("Straße wird ignoriert", "road5_test nicht in ROADS-Liste",             "Unbekannte Straße direkt verwenden",           GREEN),
    ("Training instabil",     "train/std steigt auf 3+ → Policy explodiert","Weniger Straßen, kleinere learning_rate",       YELLOW),
    ("TensorBoard leer",      "Relativer Pfad funktioniert nicht",           "Absoluten Pfad mit --logdir verwenden",        GREEN),
    ("t_end zu klein",        "config.json t_end=32 statt 320",              "Auf 320.0 zurückgesetzt",                      GREEN),
]

for i, (prob, ursache, lösung, col) in enumerate(problems):
    top = 1.7 + i * 0.67
    bg_c = DARK_GRAY if i % 2 == 0 else RGBColor(0x28, 0x28, 0x40)
    rect = s.shapes.add_shape(1, Inches(0.3), Inches(top),
                               Inches(12.7), Inches(0.6))
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_c
    rect.line.fill.background()
    box(s, 0.4,  top+0.07, 2.8, 0.5, prob,    font_size=12, bold=True, color=RED)
    box(s, 3.2,  top+0.07, 4.2, 0.5, ursache, font_size=11, color=GRAY)
    box(s, 7.4,  top+0.07, 5.4, 0.5, lösung,  font_size=11, color=col)


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 13 — Ausblick / Nächste Schritte
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
title_box(s, "Ausblick — Nächste Schritte", "Was kommt als nächstes?")
divider(s)

bullet_box(s, 0.4, 1.8, 5.8, 4.8, [
    "",
    "  Neues Training starten",
    "  Nur 1 Straße (newRoad)",
    "  learning_rate = 1e-4",
    "  ent_coef = 0.005",
    "",
    "  Modell auf road5_test testen",
    "  (Generalisierbarkeit prüfen)",
    "",
    "  TensorBoard beobachten",
    "  ep_rew_mean + ep_len_mean",
    "",
], title="  Kurzfristig:", title_color=GREEN, bg_color=DARK_GRAY)

bullet_box(s, 6.5, 1.8, 6.4, 4.8, [
    "",
    "  Gas + Bremse als Aktion",
    "  3D Action Space",
    "  [Lenkwinkel, Gas, Bremse]",
    "",
    "  Zielgeschwindigkeit als Reward",
    "  r -= 0.5 * (v - v_target)²",
    "",
    "  Vergleich: 1D vs 3D Agent",
    "  → Guter Bachelorarbeit-Inhalt!",
    "",
    "  Mehr komplexe Teststrecken",
    "",
], title="  Bachelorarbeit-Erweiterung:", title_color=YELLOW, bg_color=DARK_GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# FOLIE 14 — Zusammenfassung
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
title_box(s, "Zusammenfassung", "Was wurde in diesem Projekt erreicht?")
divider(s)

achieved = [
    "✓  PPO-Agent trainiert Spurhalten auf 4 verschiedenen Strecken",
    "✓  Zwei-Prozess-Architektur (Python 3.12 ↔ 3.13) vollständig implementiert",
    "✓  8D Observation Space (e_y, e_psi, curvature, distLeft, distRight, laneWidth, v, psidot)",
    "✓  TensorBoard-Integration für Live-Monitoring des Trainings",
    "✓  Automatische Run-Ordner (run_1, run_2, ...) — kein Überschreiben",
    "✓  vecnormalize.pkl wird nach jedem Update gespeichert",
    "✓  eval.py mit interaktiver Straßenauswahl (1–4)",
    "✓  modell_test.py für unbekannte Teststrecke (Generalisierungstest)",
    "✓  road5_test: ~1750m geschlossener Kurs mit Chicanes & S-Kurven",
    "✓  Episode endet automatisch am Streckenende",
]

for i, text in enumerate(achieved):
    top = 1.7 + i * 0.55
    box(s, 0.5, top, 12.3, 0.5, text, font_size=15,
        color=GREEN if "✓" in text else WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SPEICHERN
# ─────────────────────────────────────────────────────────────────────────────
out = "DriveLab_Praesentation.pptx"
prs.save(out)
print(f"\nPraesentation gespeichert: {out}")
print(f"  Folien: {len(prs.slides)}")
print(f"  Oeffnen mit: PowerPoint oder LibreOffice Impress")
